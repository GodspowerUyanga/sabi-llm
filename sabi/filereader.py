"""Universal file reader.

Extracts text from many formats so SABI can read, summarize and work on real
documents — PDF, Word, Excel, PowerPoint, CSV, HTML, JSON, images (OCR if
available) and any plain-text/code file. Each format degrades gracefully: if an
optional library is missing, the reader returns a clear message telling the user
which extra to install, instead of crashing.

Install the document extras with:  pip install "sabi-llm[docs]"
"""

from __future__ import annotations

import csv
import io
import json
from pathlib import Path
from typing import Optional

DEFAULT_MAX_CHARS = 6000

TEXT_EXTS = {
    ".txt", ".md", ".rst", ".log", ".py", ".js", ".ts", ".tsx", ".jsx",
    ".java", ".c", ".h", ".cpp", ".hpp", ".cc", ".go", ".rs", ".rb", ".php",
    ".sh", ".bash", ".zsh", ".ps1", ".yaml", ".yml", ".toml", ".ini", ".cfg",
    ".conf", ".env", ".sql", ".css", ".scss", ".less", ".xml", ".svg", ".tex",
    ".r", ".jl", ".kt", ".swift", ".m", ".pl", ".lua", ".dart", ".vue", ".gradle",
}
IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"}


def _truncate(text: str, max_chars: int) -> str:
    text = text.strip()
    if len(text) > max_chars:
        return text[:max_chars] + f"\n\n…[truncated, {len(text):,} chars total]"
    return text


def _need(pkg: str, ext: str) -> str:
    return (f"Cannot read {ext} files: '{pkg}' is not installed. "
            f'Install document support with:  pip install "sabi-llm[docs]"')


# --------------------------------------------------------------- per-format
def _read_pdf(p: Path) -> str:
    try:
        from pypdf import PdfReader
    except Exception:
        return _need("pypdf", ".pdf")
    reader = PdfReader(str(p))
    parts = []
    for i, page in enumerate(reader.pages, 1):
        try:
            t = page.extract_text() or ""
        except Exception:
            t = ""
        if t.strip():
            parts.append(f"[page {i}]\n{t.strip()}")
    if not parts:
        return (f"[{p.name}] is a PDF with {len(reader.pages)} page(s) but no "
                "extractable text (it may be scanned images). OCR would be needed.")
    return "\n\n".join(parts)


def _read_docx(p: Path) -> str:
    try:
        import docx
    except Exception:
        return _need("python-docx", ".docx")
    doc = docx.Document(str(p))
    out = [para.text for para in doc.paragraphs if para.text.strip()]
    for tbl in doc.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                out.append(" | ".join(cells))
    return "\n".join(out)


def _read_xlsx(p: Path) -> str:
    try:
        import openpyxl
    except Exception:
        return _need("openpyxl", ".xlsx")
    wb = openpyxl.load_workbook(str(p), read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        out.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            if any(cells):
                out.append(", ".join(cells))
    wb.close()
    return "\n".join(out)


def _read_csv(p: Path, delim: str = ",") -> str:
    out = []
    with open(p, newline="", encoding="utf-8", errors="replace") as fh:
        for row in csv.reader(fh, delimiter=delim):
            out.append(", ".join(row))
    return "\n".join(out)


def _read_pptx(p: Path) -> str:
    try:
        from pptx import Presentation
    except Exception:
        return _need("python-pptx", ".pptx")
    prs = Presentation(str(p))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"[slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                out.append(shape.text_frame.text.strip())
    return "\n".join(out)


def _read_html(p: Path) -> str:
    raw = p.read_text(encoding="utf-8", errors="replace")
    try:
        from bs4 import BeautifulSoup
    except Exception:
        return raw  # fall back to raw HTML
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text("\n", strip=True)


def _read_json(p: Path) -> str:
    try:
        data = json.loads(p.read_text(encoding="utf-8", errors="replace"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception:
        return p.read_text(encoding="utf-8", errors="replace")


def _read_image(p: Path) -> str:
    info = ""
    try:
        from PIL import Image
        with Image.open(p) as im:
            info = f"[{p.name}] image, {im.format}, {im.width}x{im.height}px."
    except Exception:
        info = f"[{p.name}] image file."
    try:
        import pytesseract
        from PIL import Image
        text = pytesseract.image_to_string(Image.open(p)).strip()
        if text:
            return info + "\n\nText found via OCR:\n" + text
        return info + "\n(no text detected via OCR)"
    except Exception:
        return info + ("\n(OCR not available — install pytesseract and the "
                       "tesseract engine to extract text from images)")


def _read_text(p: Path) -> str:
    return p.read_text(encoding="utf-8", errors="replace")


# --------------------------------------------------------------- dispatch
def read_any(path, max_chars: int = DEFAULT_MAX_CHARS) -> str:
    """Extract text from a file of (almost) any format. Always returns a string."""
    p = Path(path)
    if not p.exists():
        return f"File not found: {p}"
    if p.is_dir():
        return f"{p} is a directory, not a file."

    ext = p.suffix.lower()
    try:
        if ext == ".pdf":
            text = _read_pdf(p)
        elif ext in (".docx", ".docm"):
            text = _read_docx(p)
        elif ext in (".xlsx", ".xlsm", ".xltx"):
            text = _read_xlsx(p)
        elif ext == ".csv":
            text = _read_csv(p, ",")
        elif ext in (".tsv", ".tab"):
            text = _read_csv(p, "\t")
        elif ext in (".pptx", ".pptm"):
            text = _read_pptx(p)
        elif ext in (".html", ".htm"):
            text = _read_html(p)
        elif ext == ".json":
            text = _read_json(p)
        elif ext in IMAGE_EXTS:
            text = _read_image(p)
        elif ext in TEXT_EXTS or ext == "":
            text = _read_text(p)
        else:
            # try as text; if it looks binary, say so
            try:
                raw = p.read_bytes()[:4096]
                if b"\x00" in raw:
                    return (f"[{p.name}] appears to be a binary {ext or 'file'} "
                            "that isn't a supported document format.")
                text = p.read_text(encoding="utf-8", errors="replace")
            except Exception as exc:  # noqa: BLE001
                return f"Could not read {p.name}: {exc}"
    except Exception as exc:  # noqa: BLE001
        return f"Could not read {p.name} ({ext}): {type(exc).__name__}: {exc}"

    return _truncate(text or f"[{p.name}] contained no extractable text.", max_chars)


def supported_note() -> str:
    return ("PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx), CSV/TSV, HTML, "
            "JSON, images (OCR if available), and any text/code file.")
